# -*- coding: utf-8 -*-
"""
core/ppt.py — native, editable PowerPoint (real text, tables and charts).

Built directly from the computed Metrics (not screenshots), so slides are proper
editable objects: no image artifacts, empty slides or overlaps. 16:9.
"""
from __future__ import annotations

import config
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# palette
NAVY = RGBColor(0x06, 0x2E, 0x5C)
NAVY2 = RGBColor(0x0B, 0x3D, 0x6B)
INK = RGBColor(0x15, 0x20, 0x2B)
GREY = RGBColor(0x6B, 0x78, 0x84)
SUBTOT = RGBColor(0xEE, 0xF2, 0xF6)
EBIT_BG = RGBColor(0xFF, 0xF0, 0xF5)
ADVERSE = RGBColor(0xFC, 0xE4, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
PINK = RGBColor(0xC0, 0x39, 0x6B)
BLUE = RGBColor(0x21, 0x76, 0xAE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MINUS = "−"
EURO = "€"


# ── formatting (mirror the HTML) ─────────────────────────────────────────────
def _m1(v):
    mv = round(v / 1e6, 1)
    return 0.0 if mv == 0 else mv


def fm(v):
    mv = _m1(v)
    return f"{MINUS if mv < 0 else ''}{EURO}{abs(mv):.1f}".replace(".", ",")


def fv(v):
    mv = _m1(v)
    s = "+" if mv > 0 else (MINUS if mv < 0 else "")
    return f"{s}{EURO}{abs(mv):.1f}".replace(".", ",")


def fp(p):
    s = "+" if p > 0 else (MINUS if p < 0 else "")
    return f"{s}{abs(p):.1f}".replace(".", ",") + "%"


def _vcolor(v, fav):
    mv = _m1(v)
    if mv == 0:
        return None
    good = (mv > 0) if fav == "rev" else (mv < 0)
    return GREEN if good else PINK


# ── slide scaffolding ────────────────────────────────────────────────────────
def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bar(slide, title, sub=""):
    box = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf = box.text_frame; tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0xB6, 0xD4, 0xF0)


def _txt(slide, text, x, y, w, h, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def _table(slide, headers, rows, x, y, w, col_w=None, fs=10.5, header_fs=10.5,
           row_h=Inches(0.26)):
    """rows: list of (cells, style) where cells is list of (text, color|None, bold);
    style: 'detail'|'subtotal'|'ebitda'|'total'|'adverse'."""
    nr, nc = len(rows) + 1, len(headers)
    h = row_h * nr
    gtable = slide.shapes.add_table(nr, nc, x, y, w, h).table
    gtable.first_row = False; gtable.horz_banding = False
    if col_w:
        for i, cw in enumerate(col_w):
            gtable.columns[i].width = cw
    # header
    for c, htext in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY2
        cell.margin_left = cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext
        r.font.size = Pt(header_fs); r.font.bold = True; r.font.color.rgb = WHITE
    # body
    for ri, (cells, style) in enumerate(rows, start=1):
        bg = {"subtotal": SUBTOT, "ebitda": EBIT_BG, "total": SUBTOT,
              "adverse": ADVERSE}.get(style, WHITE)
        for c, cell_spec in enumerate(cells):
            text, color, bold = cell_spec
            cell = gtable.cell(ri, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = text
            r.font.size = Pt(fs)
            r.font.bold = bold or style in ("subtotal", "ebitda", "total")
            r.font.color.rgb = color or (PINK if style == "ebitda" and c > 0 else INK)
    return gtable


def _six_col(rows, fav_of):
    """Build (cells,style) for a 6-col results/variance table from Row objects."""
    out = []
    for r in rows:
        fav = fav_of(r)
        style = {"subtotal": "subtotal", "ebitda": "ebitda", "total": "total"}.get(
            getattr(r, "kind", "detail"), "detail")
        cells = [(r.label, None, False),
                 (fm(r.v2026), None, False), (fm(r.v2025), None, False),
                 (fm(r.budget), None, False),
                 (fv(r.vpy), _vcolor(r.vpy, fav), False),
                 (fv(r.v25), _vcolor(r.v25, fav), False)]
        out.append((cells, style))
    return out


# ── slides ───────────────────────────────────────────────────────────────────
def _hdr6(caps_month):
    return ["Δείκτης", "2026", "2025", "Π/Υ", "Απόκλ. Π/Υ", "Απόκλ. '25"]


def _title_slide(prs, m, period):
    s = _slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    _txt(s, "ΟΚΥπΥ · EBITDA Αποτελέσματα", Inches(0.7), Inches(1.2), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)
    _txt(s, period + " · Παρουσίαση Διοικητικού Συμβουλίου", Inches(0.7), Inches(2.0),
         Inches(12), Inches(0.5), size=16, color=RGBColor(0x9E, 0xCF, 0xEE))
    h = m.headline
    _txt(s, fm(h["ebitda_ytd"]) + "Μ", Inches(0.7), Inches(3.0), Inches(6), Inches(1.2),
         size=54, bold=True, color=RGBColor(0xF4, 0x8F, 0xB1))
    _txt(s, "EBITDA Περιόδου", Inches(0.75), Inches(4.2), Inches(6), Inches(0.4),
         size=14, color=WHITE)
    lines = [
        f"Π/Υ: {fm(h['ebitda_bud'])}Μ   |   Απόκλιση: {fv(h['ebitda_vpy'])}Μ",
        f"2025 (διορθ.): {fm(h['ebitda_2025'])}Μ   |   YoY: {fv(h['ebitda_yoy'])}Μ",
        f"Έσοδα {fm(h['rev_ytd'])}Μ vs Π/Υ {fm(h['rev_bud'])}Μ ({fp(h['rev_vpy_pct'])})",
        f"OpEx {fm(h['exp_ytd'])}Μ vs Π/Υ {fm(h['exp_bud'])}Μ ({fp(h['exp_vpy_pct'])})",
    ]
    _txt(s, "\n".join(lines), Inches(7.0), Inches(3.0), Inches(5.6), Inches(2.5),
         size=15, color=WHITE)


def _results_slide(prs, m, period, mgen):
    s = _slide(prs)
    _bar(s, "Σύνοψη Αποτελεσμάτων", period)
    cw = [Inches(1.9)] + [Inches(0.98)] * 5
    w = Inches(6.78)
    _txt(s, f"ΑΠΟΤΕΛΕΣΜΑΤΑ {mgen}", Inches(0.35), Inches(1.05), Inches(6), Inches(0.3),
         size=12, bold=True, color=NAVY2)
    _table(s, _hdr6(mgen), _six_col(m.results_month, lambda r: r.fav),
           Inches(0.35), Inches(1.4), w, cw)
    _txt(s, "ΣΩΡΕΥΤΙΚΑ (YTD)", Inches(0.35), Inches(4.35), Inches(6), Inches(0.3),
         size=12, bold=True, color=NAVY2)
    _table(s, _hdr6(mgen), _six_col(m.results_ytd, lambda r: r.fav),
           Inches(0.35), Inches(4.7), w, cw)
    # native clustered column: Έσοδα/Έξοδα by 2026/2025/Π/Υ
    h = m.headline
    cd = CategoryChartData(); cd.categories = ["2026", "2025", "Π/Υ"]
    cd.add_series("Έσοδα", (h["rev_ytd"] / 1e6, h["rev_2025"] / 1e6, h["rev_bud"] / 1e6))
    cd.add_series("Έξοδα", (h["exp_ytd"] / 1e6, h["exp_2025"] / 1e6, h["exp_bud"] / 1e6))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.4), Inches(1.4),
                            Inches(5.5), Inches(3.0), cd).chart
    gf.has_title = True; gf.chart_title.text_frame.text = "Έσοδα vs Έξοδα (€Μ)"
    gf.has_legend = True; gf.legend.position = XL_LEGEND_POSITION.BOTTOM; gf.legend.include_in_layout = False
    # EBITDA by month line
    mo = m.monthly
    labels = [config.MONTHS[k + 1]["short"] for k in range(m.mm)]
    cd2 = CategoryChartData(); cd2.categories = labels
    cd2.add_series("EBITDA 2026", tuple(x / 1e6 for x in mo["ebitda26"]))
    cd2.add_series("EBITDA 2025", tuple(x / 1e6 for x in mo["ebitda25"]))
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(7.4), Inches(4.5),
                             Inches(5.5), Inches(2.6), cd2).chart
    gf2.has_title = True; gf2.chart_title.text_frame.text = "EBITDA ανά μήνα (€Μ)"
    gf2.has_legend = True; gf2.legend.position = XL_LEGEND_POSITION.BOTTOM; gf2.legend.include_in_layout = False


def _pl_slide(prs, m, period):
    s = _slide(prs)
    _bar(s, "Αναλυτικές Αποτελεσμάτων — Ενοποιημένο EBITDA P&L", period)
    ov = m.overview
    rows = []
    for c in ov["rev_cats"]:
        vpy = c["ytd"] - c["bud"]; yoy = c["ytd"] - c["y25"]
        rows.append(([(c["l"], None, False), (fm(c["ytd"]), None, False),
                      (fm(c["bud"]), None, False), (fv(vpy), _vcolor(vpy, "rev"), False),
                      (fm(c["y25"]), None, False), (fv(yoy), _vcolor(yoy, "rev"), False)], "detail"))
    h = m.headline
    rows.append(([("ΣΥΝΟΛΟ ΕΣΟΔΩΝ", None, True), (fm(h["rev_ytd"]), None, True),
                  (fm(h["rev_bud"]), None, True), (fv(h["rev_vpy"]), _vcolor(h["rev_vpy"], "rev"), True),
                  (fm(h["rev_2025"]), None, True), (fv(h["rev_ytd"] - h["rev_2025"]), None, True)], "subtotal"))
    for c in ov["exp_cats"]:
        vpy = c["ytd"] - c["bud"]; yoy = c["ytd"] - c["y25"]
        rows.append(([(c["l"], None, False), (fm(c["ytd"]), None, False),
                      (fm(c["bud"]), None, False), (fv(vpy), _vcolor(vpy, "exp"), False),
                      (fm(c["y25"]), None, False), (fv(yoy), _vcolor(yoy, "exp"), False)], "detail"))
    rows.append(([("ΣΥΝΟΛΟ OPEX", None, True), (fm(h["exp_ytd"]), None, True),
                  (fm(h["exp_bud"]), None, True), (fv(h["exp_vpy"]), _vcolor(h["exp_vpy"], "exp"), True),
                  (fm(h["exp_2025"]), None, True), (fv(h["exp_ytd"] - h["exp_2025"]), None, True)], "subtotal"))
    rows.append(([("EBITDA", None, True), (fm(h["ebitda_ytd"]), None, True),
                  (fm(h["ebitda_bud"]), None, True), (fv(h["ebitda_vpy"]), None, True),
                  (fm(h["ebitda_2025"]), None, True), (fv(h["ebitda_yoy"]), None, True)], "ebitda"))
    cw = [Inches(3.6)] + [Inches(1.7)] * 5
    _table(s, ["Κατηγορία", "YTD 2026", "Π/Υ", "Απόκλ. Π/Υ", "YTD 2025", "YoY"],
           rows, Inches(0.5), Inches(1.15), Inches(12.1), cw, fs=9.5, row_h=Inches(0.23))


def _simple_table_slide(prs, title, period, headers, rows, cw, x=Inches(0.5),
                        w=Inches(12.3), fs=10, row_h=Inches(0.28)):
    s = _slide(prs)
    _bar(s, title, period)
    _table(s, headers, rows, x, Inches(1.2), w, cw, fs=fs, row_h=row_h)
    return s


def _hosp_slide(prs, m, period):
    rows = []
    for h in m.hospitals:
        n26 = h["r"] - h["e"]; rv = (h["r"] - h["rb"]) / h["rb"] * 100 if h["rb"] else 0
        ev = (h["e"] - h["eb"]) / h["eb"] * 100 if h["eb"] else 0
        rows.append(([(h["n"], None, False), (fm(h["r"]), None, False),
                      (fp(rv), _vcolor(h["r"] - h["rb"], "rev"), False),
                      (fm(h["e"]), None, False), (fp(ev), _vcolor(h["e"] - h["eb"], "exp"), False),
                      (fv(n26), _vcolor(n26, "rev"), False),
                      (fv(h["n25"]), _vcolor(h["n25"], "rev"), False)], "detail"))
    tr = sum(h["r"] for h in m.hospitals); te = sum(h["e"] for h in m.hospitals)
    rows.append(([("ΣΥΝΟΛΟ", None, True), (fm(tr), None, True), ("", None, True),
                  (fm(te), None, True), ("", None, True), (fv(tr - te), None, True),
                  (fv(m.headline["ebitda_2025"]), None, True)], "ebitda"))
    cw = [Inches(2.7), Inches(1.5), Inches(1.3), Inches(1.5), Inches(1.3), Inches(1.6), Inches(1.6)]
    s = _simple_table_slide(prs, "Ανά Νοσηλευτήριο", period,
                            ["Νοσηλευτήριο", "Έσοδα", "Έσ.%", "OpEx", "Εξ.%", "Καθαρό '26", "Καθαρό '25"],
                            rows, cw, fs=9.5, row_h=Inches(0.30))
    # native bar: rev vs opex by hospital
    cd = CategoryChartData(); cd.categories = [h["n"] for h in m.hospitals]
    cd.add_series("Έσοδα", tuple(h["r"] / 1e6 for h in m.hospitals))
    cd.add_series("OpEx", tuple(h["e"] / 1e6 for h in m.hospitals))
    # (chart omitted here to keep the table full-width & readable)


def _oay_slide(prs, m, period):
    rows = []
    for c in m.oay["cats"]:
        vpy = c["ytd"] - c["bud"]; yoy = c["ytd"] - c["y25"]
        rows.append(([(c["label"], None, False), (fm(c["ytd"]), None, False),
                      (fm(c["bud"]), None, False), (fv(vpy), _vcolor(vpy, "rev"), False),
                      (fm(c["y25"]), None, False), (fv(yoy), _vcolor(yoy, "rev"), False)], "detail"))
    t = m.oay["total"]
    rows.append(([("ΣΥΝΟΛΟ ΟΑΥ", None, True), (fm(t["ytd"]), None, True),
                  (fm(t["bud"]), None, True), (fv(t["ytd"] - t["bud"]), _vcolor(t["ytd"] - t["bud"], "rev"), True),
                  (fm(t["y25"]), None, True), (fv(t["ytd"] - t["y25"]), None, True)], "subtotal"))
    cw = [Inches(4.0)] + [Inches(1.66)] * 5
    _simple_table_slide(prs, "ΟΑΥ Έσοδα", period,
                        ["Κατηγορία", "YTD 2026", "Π/Υ", "Απόκλ. Π/Υ", "YTD 2025", "YoY"], rows, cw)


def _alla_slide(prs, m, period):
    rows = []
    for c in m.allaesoda["subs"][:12]:
        vpy = c["ytd"] - c["bud"]
        rows.append(([(c["label"], None, False), (fm(c["ytd"]), None, False),
                      (fm(c["bud"]), None, False), (fv(vpy), _vcolor(vpy, "rev"), False),
                      (fm(c["y25"]), None, False)], "detail"))
    t = m.allaesoda["total"]
    rows.append(([("ΣΥΝΟΛΟ", None, True), (fm(t["ytd"]), None, True), (fm(t["bud"]), None, True),
                  (fv(t["ytd"] - t["bud"]), _vcolor(t["ytd"] - t["bud"], "rev"), True),
                  (fm(t["y25"]), None, True)], "subtotal"))
    cw = [Inches(4.8), Inches(1.9), Inches(1.9), Inches(1.9), Inches(1.8)]
    _simple_table_slide(prs, "Άλλα Έσοδα", period,
                        ["Υποκατηγορία", "YTD 2026", "Π/Υ", "Απόκλ.", "YTD 2025"], rows, cw,
                        fs=9.5, row_h=Inches(0.30))


def _loipa_slide(prs, m, period):
    lp = m.loipa; rows = []
    for c in lp["cats"]:
        vpy = c["ytd"] - c["bud"]
        rows.append(([(c["label"], None, False), (fm(c["ytd"]), None, False),
                      (fm(c["y25"]), None, False), (fm(c["bud"]), None, False),
                      (fv(vpy), _vcolor(vpy, "exp"), False)], "detail"))
    t = lp["total"]
    rows.append(([("ΣΥΝΟΛΟ Λειτ. Εξόδων", None, True), (fm(t["ytd"]), None, True),
                  (fm(t["y25"]), None, True), (fm(t["bud"]), None, True),
                  (fv(t["ytd"] - t["bud"]), _vcolor(t["ytd"] - t["bud"], "exp"), True)], "subtotal"))
    cw = [Inches(4.8), Inches(1.9), Inches(1.9), Inches(1.9), Inches(1.8)]
    _simple_table_slide(prs, "Λειτουργικά Έξοδα (εκτός Μισθοδοσίας)", period,
                        ["Κατηγορία", "YTD 2026", "2025", "Π/Υ", "Απόκλ. Π/Υ"], rows, cw)


def _payroll_slide(prs, m, period):
    s = _slide(prs)
    _bar(s, "Μισθοδοσία & Κόστος Προσωπικού", period)
    p = m.payroll
    rows = []
    for r in p["summary"]:
        vpy = r["ytd"] - r["bud"]; yoy = r["ytd"] - r["y25"]
        rows.append(([(r["label"], None, False), (fm(r["ytd"]), None, False),
                      (fm(r["bud"]), None, False), (fv(vpy), _vcolor(vpy, "exp"), False),
                      (fm(r["y25"]), None, False), (fv(yoy), _vcolor(yoy, "exp"), False)], "detail"))
    t = p["total"]
    rows.append(([("ΣΥΝΟΛΟ ΠΡΟΣΩΠΙΚΟΥ", None, True), (fm(t["ytd"]), None, True),
                  (fm(t["bud"]), None, True), (fv(t["ytd"] - t["bud"]), None, True),
                  (fm(t["y25"]), None, True), (fv(t["ytd"] - t["y25"]), None, True)], "ebitda"))
    _table(s, ["Κατηγορία", "YTD 2026", "Π/Υ", "Απόκλ.", "2025", "YoY"], rows,
           Inches(0.5), Inches(1.2), Inches(7.0),
           [Inches(2.4)] + [Inches(0.92)] * 5, fs=10, row_h=Inches(0.3))
    # Σύμβαση components (compact)
    comp = []
    for c in p["symv_components"]:
        vpy = c["ytd"] - c["bud"]
        comp.append(([(c["label"], None, False), (fm(c["ytd"]), None, False),
                      (fm(c["bud"]), None, False), (fv(vpy), _vcolor(vpy, "exp"), False)], "detail"))
    _txt(s, "Σύμβαση ΟΚΥπΥ — κατηγορίες", Inches(7.8), Inches(1.2), Inches(5), Inches(0.3),
         size=12, bold=True, color=NAVY2)
    _table(s, ["Κατηγορία", "2026", "Π/Υ", "Απόκλ."], comp, Inches(7.8), Inches(1.55),
           Inches(5.0), [Inches(2.2), Inches(1.0), Inches(1.0), Inches(0.8)], fs=9.5,
           row_h=Inches(0.26))


def build_pptx(m, out_path: str, mm: int, year: int) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    gmon = config.MONTHS[mm]
    period = f"Ιαν–{gmon['caps_short']} {year}"
    mgen = f"{gmon['caps_gen']} {year}"
    _title_slide(prs, m, period)
    _results_slide(prs, m, period, mgen)
    _pl_slide(prs, m, period)
    _hosp_slide(prs, m, period)
    _oay_slide(prs, m, period)
    _alla_slide(prs, m, period)
    _loipa_slide(prs, m, period)
    _payroll_slide(prs, m, period)
    prs.save(out_path)
    return out_path
